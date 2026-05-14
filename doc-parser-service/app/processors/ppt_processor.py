import time
import uuid
from pathlib import Path
from loguru import logger
from pptx import Presentation


from app.core.base_parser import BaseParser
from app.models.response import (
    ParseResult,
    DocumentMetadata,
    LogicalSegment,
    PageDebugInfo,
    LayoutBlock,
    BlockType,
    BoundingBox,
)


class PptProcessor(BaseParser):
    def supports(self, file_extension: str, strategy: str) -> bool:
        return file_extension in ("pptx", "ppt")

    def order(self) -> int:
        return 10

    def parse(self, file_path: Path, strategy: str, language: str, **kwargs) -> ParseResult:
        start = time.time()
        try:
            prs = Presentation(str(file_path))

            segments: list[LogicalSegment] = []
            pages: list[PageDebugInfo] = []
            order = 0

            for slide_idx, slide in enumerate(prs.slides):
                slide_blocks: list[LayoutBlock] = []

                for shape in slide.shapes:
                    block_id = str(uuid.uuid4())

                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if not text:
                            continue

                        is_title = shape.shape_type is not None and "title" in str(shape.name).lower()
                        block_type = BlockType.TITLE if is_title else BlockType.TEXT

                        slide_blocks.append(LayoutBlock(
                            id=block_id,
                            type=block_type,
                            bbox=BoundingBox(
                                x0=float(shape.left),
                                y0=float(shape.top),
                                x1=float(shape.left + shape.width),
                                y1=float(shape.top + shape.height),
                            ),
                            text=text,
                            page_num=slide_idx + 1,
                            confidence=1.0,
                        ))

                        segments.append(LogicalSegment(
                            id=block_id,
                            type="heading" if is_title else "paragraph",
                            content=text,
                            start_order=order,
                            end_order=order,
                        ))
                        order += 1

                    elif shape.has_table:
                        table_html = self._table_to_html(shape.table)

                        slide_blocks.append(LayoutBlock(
                            id=block_id,
                            type=BlockType.TABLE,
                            bbox=BoundingBox(
                                x0=float(shape.left),
                                y0=float(shape.top),
                                x1=float(shape.left + shape.width),
                                y1=float(shape.top + shape.height),
                            ),
                            html=table_html,
                            page_num=slide_idx + 1,
                            confidence=1.0,
                        ))

                        segments.append(LogicalSegment(
                            id=block_id,
                            type="table",
                            content=table_html,
                            start_order=order,
                            end_order=order,
                        ))
                        order += 1

                pages.append(PageDebugInfo(
                    page_num=slide_idx + 1,
                    column_count=1,
                    blocks=slide_blocks,
                ))

            elapsed = int((time.time() - start) * 1000)
            return ParseResult(
                success=True,
                segments=segments,
                pages=pages,
                document_metadata=DocumentMetadata(
                    filename=file_path.name,
                    file_type=file_path.suffix.lstrip("."),
                    page_count=len(prs.slides),
                    language=language,
                    parse_strategy=strategy,
                    parser_engine="python-pptx",
                ),
                processing_time_ms=elapsed,
            )
        except Exception as e:
            logger.exception("PowerPoint parse failed")
            elapsed = int((time.time() - start) * 1000)
            return ParseResult(success=False, error_message=str(e), processing_time_ms=elapsed)

    def _table_to_html(self, table) -> str:
        html = "<table>"
        for row in table.rows:
            html += "<tr>"
            for cell in row.cells:
                text = cell.text.strip()
                html += f"<td>{text}</td>"
            html += "</tr>"
        html += "</table>"
        return html
